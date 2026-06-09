from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

def create_polished_ppt(template_path, output_path, content_data):
    """
    Args:
        template_path (str): 你们的高端设计感 PPT 模板路径 (.pptx)
        output_path (str): 生成的规范化 PPT 路径
        content_data (list): 润色后的结构化数据，例如: 
                             [{"layout_idx": 0, "title": "封面", "body": "副标题"}, ...]
    """
    # 1. 加载你提供的品牌模板
    prs = Presentation(template_path)
    
    # 2. 遍历数据，严格根据模板布局生成页面
    for data in content_data:
        layout_idx = data.get("layout_idx", 1) # 默认使用内容页布局
        
        # 严格调用模板自带的母版布局，确保背景、LOGO、字体样式纯正
        slide_layout = prs.slide_layouts[layout_idx]
        slide = prs.slides.add_slide(slide_layout)
        
        # 3. 精准填入模板预设的占位符 (Placeholders)，绝不乱新建文本框
        # 通常 placeholders[0] 是标题，placeholders[1] 是正文/副标题
        if slide.shapes.title and "title" in data:
            slide.shapes.title.text = data["title"]
            
        if len(slide.placeholders) > 1 and "body" in data:
            body_shape = slide.placeholders[1]
            tf = body_shape.text_frame
            
            # 如果是列表数据，自动转为高级的 Bullet Points 或分段
            if isinstance(data["body"], list):
                tf.text = data["body"][0]
                for bullet in data["body"][1:]:
                    p = tf.add_paragraph()
                    p.text = bullet
                    p.level = 1 # 保持缩进
            else:
                tf.text = data["body"]
                
    # 4. 保存文件
    prs.save(output_path)
    print(f"✨ 严格遵循模板规范的 PPT 已成功输出至: {output_path}")

# ==========================================
# Manus 执行时的模拟测试数据
# ==========================================
if __name__ == "__main__":
    # 假设用户上传了模板 template.pptx
    # layout_idx 映射：0-封面, 1-纯文字内容, 2-图文版式（具体依你的模板而定）
    mock_polished_data = [
        {
            "layout_idx": 0, 
            "title": "2026 公共休闲办公空间趋势", 
            "body": "基于极简美学与人体工学的软体家具重塑"
        },
        {
            "layout_idx": 1, 
            "title": "设计核心：有价值感的留白", 
            "body": [
                "延续 Arper 式的视觉呼吸感，拒绝高饱和度色块堆砌。",
                "模块化沙发系统：通过流畅的线条勾勒非正式协作空间。",
                "面料与质感：采用低调的高级中性色，凸显公共空间的温暖包裹感。"
            ]
        }
    ]
    
    # 在 Manus 中运行时，它会自动替换为实际的路径
    try:
        create_polished_ppt("brand_template.pptx", "output_polished_presentation.pptx", mock_polished_data)
    except Exception as e:
        print(f"执行失败，请确保沙箱中存在 brand_template.pptx 模板。错误: {e}")
